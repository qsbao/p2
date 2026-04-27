"""Structural extractor (Q2, Q5, Q6).

Per slide, writes:
  <debug_dir>/manifest-{N}.json
  <debug_dir>/notes-{N}.txt
  <debug_dir>/chrome-{N}.json
  <media_dir>/slide{N}-fig{M}.png             (cropped from rendered slide PNG)
  <media_dir>/slide{N}-fig{M}.original.<ext>  (only for <p:pic>)

The manifest is the ground-truth artifact the VLM consumes. Every text string the
VLM emits must be a substring of some manifest shape's `text` field.
"""

from __future__ import annotations

import json
from dataclasses import asdict, dataclass, field
from pathlib import Path
from typing import Any

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.group import GroupShape


# Group bbox must cover ≥5% of slide area to count as image-like (Q5).
GROUP_AREA_THRESHOLD = 0.05
# Pad crops by 2% of slide width so anti-aliased edges aren't clipped (Q5).
CROP_PAD_FRAC = 0.02


@dataclass
class ShapeRecord:
    id: str
    kind: str  # "Picture" | "Group" | "GraphicFrame" | "Shape"
    autoshape: str | None  # e.g. "TEXT_BOX", "RECTANGLE"; None when not an autoshape
    name: str
    z: int
    bbox_emu: list[int]  # [left, top, width, height]
    bbox_frac: list[float]  # [left, top, width, height], normalized 0..1
    text: str  # paragraph-joined plain text; "" if no text frame
    paragraphs: list[dict[str, Any]]  # [{"runs": [{"text", "size_pt", "bold", "italic", "font"}]}]
    is_placeholder: bool
    is_master_inherited: bool
    image_path: str | None  # relative to <out_root>, e.g. "media/slide1-fig3.png"
    original_blob_path: str | None  # only set for <p:pic>


@dataclass
class SlideArtifact:
    slide_index: int
    manifest_path: Path
    notes_path: Path
    chrome_path: Path
    image_paths: list[Path] = field(default_factory=list)


def _autoshape_name(shape) -> str | None:
    """Return the MSO_SHAPE_TYPE enum name (TEXT_BOX, RECTANGLE, ...) or None."""
    try:
        st = shape.shape_type
    except Exception:
        return None
    if st is None:
        return None
    return getattr(st, "name", None)


def _is_master_inherited(shape) -> bool:
    """Conservative: a placeholder whose text frame is fully empty (i.e. content
    comes from the layout/master). At N=1 in v1 we only drop confirmed cases."""
    if not getattr(shape, "is_placeholder", False):
        return False
    if not getattr(shape, "has_text_frame", False):
        return True  # empty placeholder slot, content from master
    return shape.text_frame.text.strip() == ""


# DrawingML namespace — `<a:t>` is the universal text-run element used inside
# charts, tables, SmartArt drawings, and ordinary text frames.
_DRAWINGML_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
_A_T = f"{{{_DRAWINGML_NS}}}t"


def _graphic_frame_text(shape) -> tuple[str, list[dict[str, Any]]]:
    """Surface visible text from a GraphicFrame (chart / table / SmartArt drawing).

    python-pptx exposes `<p:graphicFrame>` as a shape but doesn't model its inner
    text runs (no `text_frame`). We walk the lxml element tree for every `<a:t>`
    descendant — that catches chart category labels, table cell text, and SmartArt
    fallback drawing text. Each run becomes its own paragraph so the substring
    validator in vlm.py can match phrase-level emissions.
    """
    el = getattr(shape, "_element", None)
    if el is None:
        return "", []
    paragraphs: list[dict[str, Any]] = []
    runs_text: list[str] = []
    for t in el.iter(_A_T):
        if t.text:
            runs_text.append(t.text)
            paragraphs.append({"runs": [{"text": t.text, "font": None, "size_pt": None, "bold": None, "italic": None}]})
    return "\n".join(runs_text), paragraphs


def _extract_paragraphs(shape) -> tuple[str, list[dict[str, Any]]]:
    """Return (joined_text, paragraph_list)."""
    if not getattr(shape, "has_text_frame", False):
        if shape.__class__.__name__ == "GraphicFrame":
            return _graphic_frame_text(shape)
        return "", []
    paragraphs = []
    for p in shape.text_frame.paragraphs:
        runs = []
        for r in p.runs:
            font = r.font
            size_pt = float(font.size.pt) if font.size is not None else None
            runs.append(
                {
                    "text": r.text,
                    "font": font.name,
                    "size_pt": size_pt,
                    "bold": bool(font.bold) if font.bold is not None else None,
                    "italic": bool(font.italic) if font.italic is not None else None,
                }
            )
        # If no runs were extracted but the paragraph has text (rare; defaults), keep it.
        if not runs and p.text:
            runs.append({"text": p.text, "font": None, "size_pt": None, "bold": None, "italic": None})
        paragraphs.append({"runs": runs})
    text = shape.text_frame.text
    return text, paragraphs


def _flatten(shape, slide_idx: int, slide_size_emu: tuple[int, int], counter: list[int]) -> list[ShapeRecord]:
    """Recurse into groups, emitting one ShapeRecord per leaf + one for each group itself.

    `counter` is a 1-element list used as a mutable z-counter.
    """
    out: list[ShapeRecord] = []
    z = counter[0]
    counter[0] += 1

    sw, sh = slide_size_emu
    left = int(shape.left or 0)
    top = int(shape.top or 0)
    width = int(shape.width or 0)
    height = int(shape.height or 0)
    bbox_emu = [left, top, width, height]
    bbox_frac = [left / sw, top / sh, width / sw, height / sh]

    is_group = isinstance(shape, GroupShape) or _autoshape_name(shape) == "GROUP"
    is_picture = _autoshape_name(shape) == "PICTURE"
    is_graphic_frame = _autoshape_name(shape) == "CHART" or _autoshape_name(shape) == "DIAGRAM" or shape.__class__.__name__ == "GraphicFrame"

    if is_group:
        kind = "Group"
    elif is_picture:
        kind = "Picture"
    elif is_graphic_frame:
        kind = "GraphicFrame"
    else:
        kind = "Shape"

    text, paragraphs = _extract_paragraphs(shape)

    rec = ShapeRecord(
        id=f"s{slide_idx}-{z}",
        kind=kind,
        autoshape=_autoshape_name(shape),
        name=shape.name,
        z=z,
        bbox_emu=bbox_emu,
        bbox_frac=bbox_frac,
        text=text,
        paragraphs=paragraphs,
        is_placeholder=bool(getattr(shape, "is_placeholder", False)),
        is_master_inherited=_is_master_inherited(shape),
        image_path=None,
        original_blob_path=None,
    )
    out.append(rec)

    if is_group:
        for child in shape.shapes:
            out.extend(_flatten(child, slide_idx, slide_size_emu, counter))

    return out


def _is_image_like(rec: ShapeRecord) -> bool:
    """Q5 image-like: <p:pic>, <p:grpSp> ≥5% slide area, <p:graphicFrame>."""
    if rec.kind == "Picture":
        return True
    if rec.kind == "GraphicFrame":
        return True
    if rec.kind == "Group":
        _, _, w_frac, h_frac = rec.bbox_frac
        return (w_frac * h_frac) >= GROUP_AREA_THRESHOLD
    return False


def _crop_image(slide_png: Path, rec: ShapeRecord, out_path: Path) -> None:
    """Crop the rec's bbox from slide_png with 2% padding, save as PNG."""
    img = Image.open(slide_png)
    W_px, H_px = img.size
    l_frac, t_frac, w_frac, h_frac = rec.bbox_frac
    pad_px = int(CROP_PAD_FRAC * W_px)
    l_px = max(0, int(l_frac * W_px) - pad_px)
    t_px = max(0, int(t_frac * H_px) - pad_px)
    r_px = min(W_px, int((l_frac + w_frac) * W_px) + pad_px)
    b_px = min(H_px, int((t_frac + h_frac) * H_px) + pad_px)
    crop = img.crop((l_px, t_px, r_px, b_px))
    out_path.parent.mkdir(parents=True, exist_ok=True)
    crop.save(out_path, "PNG")


def _save_original_blob(shape, out_path_no_ext: Path) -> Path | None:
    """For <p:pic> only: save the raw embedded image bytes. Returns the saved path."""
    try:
        image = shape.image
    except Exception:
        return None
    # `image.ext` opens the blob with PIL to sniff the format, which raises
    # UnidentifiedImageError on EMF/WMF and other vector formats PIL can't read.
    # `image.filename` is set by python-pptx to "image.<partname-ext>" when the
    # blob comes from the package — no PIL involved, so it's safe.
    try:
        ext = image.ext or "png"
    except Exception:
        fn = image.filename or ""
        ext = fn.rsplit(".", 1)[1].lower() if "." in fn else "bin"
    out = out_path_no_ext.with_suffix(f".original.{ext}")
    out.parent.mkdir(parents=True, exist_ok=True)
    out.write_bytes(image.blob)
    return out


def _picture_shapes_in_z_order(slide) -> dict[str, Any]:
    """Map ShapeRecord.id → underlying python-pptx shape for Picture-kind records.

    Needed because the manifest stores serializable data, but blob extraction needs
    the live shape object."""
    by_id: dict[str, Any] = {}
    counter = [0]

    def walk(parent_shapes):
        for sh in parent_shapes:
            z = counter[0]
            counter[0] += 1
            by_id[f"__z{z}"] = sh
            if isinstance(sh, GroupShape):
                walk(sh.shapes)

    walk(slide.shapes)
    return by_id


def _extract_notes(slide) -> str:
    if not slide.has_notes_slide:
        return ""
    return slide.notes_slide.notes_text_frame.text or ""


def extract(
    pptx: Path,
    slide_pngs: list[Path],
    out_root: Path,
    debug_dir: Path,
) -> list[SlideArtifact]:
    """Build the manifest + notes + chrome + media crops for each slide.

    out_root: directory that will hold `media/`. cli.py uses `<out_dir>/<stem>/`.
    debug_dir: directory for `manifest-{N}.json`, `notes-{N}.txt`, `chrome-{N}.json`.
    """
    pptx = Path(pptx).resolve()
    out_root = Path(out_root).resolve()
    debug_dir = Path(debug_dir).resolve()
    out_root.mkdir(parents=True, exist_ok=True)
    debug_dir.mkdir(parents=True, exist_ok=True)
    media_dir = out_root / "media"
    media_dir.mkdir(parents=True, exist_ok=True)

    prs = Presentation(pptx)
    slide_size_emu = (prs.slide_width, prs.slide_height)

    artifacts: list[SlideArtifact] = []
    for idx, slide in enumerate(prs.slides, start=1):
        slide_png = slide_pngs[idx - 1]
        slide_size_px = list(Image.open(slide_png).size)

        # Walk shapes, building records.
        records: list[ShapeRecord] = []
        z_counter = [0]
        for sh in slide.shapes:
            records.extend(_flatten(sh, idx, slide_size_emu, z_counter))

        # Map z → live shape for blob extraction.
        z_to_shape = _picture_shapes_in_z_order(slide)

        # Crop image-like shapes (and save original blob for Pictures).
        fig_idx = 0
        image_paths: list[Path] = []
        for rec in records:
            if not _is_image_like(rec):
                continue
            fig_idx += 1
            crop_name = f"slide{idx}-fig{fig_idx}.png"
            crop_path = media_dir / crop_name
            _crop_image(slide_png, rec, crop_path)
            rec.image_path = f"media/{crop_name}"
            image_paths.append(crop_path)
            if rec.kind == "Picture":
                live = z_to_shape.get(f"__z{rec.z}")
                if live is not None:
                    blob_path = _save_original_blob(live, crop_path.with_suffix(""))
                    if blob_path is not None:
                        rec.original_blob_path = f"media/{blob_path.name}"

        # Build chrome.json — at v1 only "high" (master-inherited) is computed.
        chrome = {
            "high": [r.id for r in records if r.is_master_inherited],
            "medium": [],  # cross-slide repetition deferred (N=1)
            "low": [],
        }

        # Write manifest.
        manifest = {
            "slide_index": idx,
            "slide_size_emu": list(slide_size_emu),
            "slide_size_px": slide_size_px,
            "shapes": [asdict(r) for r in records],
        }
        manifest_path = debug_dir / f"manifest-{idx}.json"
        manifest_path.write_text(json.dumps(manifest, indent=2, ensure_ascii=False))

        # Write notes.
        notes_path = debug_dir / f"notes-{idx}.txt"
        notes_path.write_text(_extract_notes(slide))

        # Write chrome.
        chrome_path = debug_dir / f"chrome-{idx}.json"
        chrome_path.write_text(json.dumps(chrome, indent=2))

        artifacts.append(
            SlideArtifact(
                slide_index=idx,
                manifest_path=manifest_path,
                notes_path=notes_path,
                chrome_path=chrome_path,
                image_paths=image_paths,
            )
        )

    return artifacts


if __name__ == "__main__":
    import sys

    from .render import render

    if len(sys.argv) != 3:
        print("usage: python -m ppt2md.extract <pptx> <out_dir>", file=sys.stderr)
        sys.exit(2)
    pptx = Path(sys.argv[1])
    out_dir = Path(sys.argv[2])
    stem = pptx.stem
    debug_dir = out_dir / f"{stem}.debug"
    out_root = out_dir / stem
    pngs = render(pptx, debug_dir)
    arts = extract(pptx, pngs, out_root, debug_dir)
    for a in arts:
        print(f"slide {a.slide_index}: manifest={a.manifest_path} notes={a.notes_path} chrome={a.chrome_path} figs={len(a.image_paths)}")
